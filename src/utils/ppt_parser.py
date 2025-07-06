import os
from pptx import Presentation
import PyPDF2

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint (.pptx) file."""
    text_content = []
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"The file {file_path} does not exist.")
    
    presentation = Presentation(file_path)
    
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    
    return text_content

def extract_text_from_pdf(file_path):
    """Extract text from PDF file."""
    text_content = []
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"The file {file_path} does not exist.")
    
    with open(file_path, "rb") as pdf_file:
        reader = PyPDF2.PdfReader(pdf_file)
        for page in reader.pages:
            text_content.append(page.extract_text())
    
    return text_content

def parse_file(file):
    """Parse the uploaded file (PPTX or PDF) and return extracted text."""
    file_path = f"/tmp/{file.name}"
    with open(file_path, "wb") as f:
        f.write(file.getbuffer())
    
    if file.name.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    elif file.name.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    else:
        raise ValueError("Unsupported file format. Only .pptx and .pdf are supported.")